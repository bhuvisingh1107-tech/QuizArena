"""Document text extraction for AI quiz generation."""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

from app.core.exceptions import ValidationError

logger = logging.getLogger(__name__)

# Formats we can extract when optional deps are installed.
ALLOWED_EXTENSIONS = {
    ".pdf",
    ".pptx",
    ".docx",
    ".png",
    ".jpg",
    ".jpeg",
    ".mp4",
    ".txt",
}

# Accepted at upload gate with clear conversion guidance (not silently broken).
LEGACY_UNSUPPORTED = {
    ".ppt": "PowerPoint 97-2003 (.ppt) is not supported. Please re-save as .pptx and upload again.",
    ".doc": "Word 97-2003 (.doc) is not supported. Please re-save as .docx and upload again.",
}

UPLOAD_ACCEPT_EXTENSIONS = ALLOWED_EXTENSIONS | set(LEGACY_UNSUPPORTED.keys())


@dataclass
class ExtractionResult:
    text: str
    extractor: str


def extract_text(path: Path, *, mime_type: str = "", filename: str = "") -> ExtractionResult:
    suffix = path.suffix.lower() or Path(filename).suffix.lower()

    if suffix in LEGACY_UNSUPPORTED:
        raise ValidationError("UNSUPPORTED_FILE_TYPE", LEGACY_UNSUPPORTED[suffix])

    try:
        if suffix == ".txt" or mime_type.startswith("text/plain"):
            return ExtractionResult(
                text=path.read_text(encoding="utf-8", errors="ignore"),
                extractor="txt",
            )
        if suffix == ".pdf" or mime_type == "application/pdf":
            return _extract_pdf(path)
        if suffix == ".pptx" or "presentation" in mime_type:
            return _extract_pptx(path)
        if suffix == ".docx" or "wordprocessingml" in mime_type:
            return _extract_docx(path)
        if suffix in {".png", ".jpg", ".jpeg"} or mime_type.startswith("image/"):
            return _extract_image_ocr(path)
        if suffix == ".mp4" or mime_type.startswith("video/"):
            return _extract_video_whisper(path)
    except ValidationError:
        raise
    except Exception as exc:
        logger.exception("Extraction failed for %s", filename or path.name)
        raise ValidationError(
            "EXTRACTION_FAILED",
            f"Could not extract text from '{filename or path.name}'. "
            "The file may be corrupt, empty, or password-protected.",
            details=[{"reason": type(exc).__name__}],
        ) from exc

    raise ValidationError(
        "UNSUPPORTED_FILE_TYPE",
        f"Unsupported file type '{suffix or mime_type}'. "
        f"Supported: {', '.join(sorted(ALLOWED_EXTENSIONS))}",
    )


def _missing_dep(code: str, user_message: str, *, cause: Exception) -> ValidationError:
    logger.error("%s: %s", code, cause)
    return ValidationError(code, user_message)


def _extract_pdf(path: Path) -> ExtractionResult:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise _missing_dep(
            "EXTRACTION_UNAVAILABLE",
            "PDF processing is temporarily unavailable. Please try a TXT or DOCX file, or contact support.",
            cause=exc,
        ) from exc

    doc = fitz.open(path)
    parts: list[str] = []
    try:
        for page_index, page in enumerate(doc):
            text = page.get_text("text") or ""
            if text.strip():
                parts.append(f"## Page {page_index + 1}\n{text.strip()}")
            # Attempt simple table-ish block extraction via dict blocks
            try:
                blocks = page.get_text("dict").get("blocks") or []
                for block in blocks:
                    if block.get("type") != 0:
                        continue
                    for line in block.get("lines") or []:
                        line_text = "".join(
                            span.get("text", "") for span in (line.get("spans") or [])
                        ).strip()
                        if line_text and line_text not in text:
                            parts.append(line_text)
            except Exception:
                pass
    finally:
        doc.close()

    combined = "\n\n".join(parts).strip()
    if not combined:
        raise ValidationError(
            "EMPTY_EXTRACTION",
            "No readable text was found in this PDF. It may be image-only — try OCR via image upload or a text-based PDF.",
        )
    return ExtractionResult(text=combined, extractor="pymupdf")


def _extract_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise _missing_dep(
            "EXTRACTION_UNAVAILABLE",
            "PowerPoint processing is temporarily unavailable. Please try PDF, DOCX, or TXT.",
            cause=exc,
        ) from exc

    prs = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_bits: list[str] = [f"## Slide {index}"]
        title = ""
        if slide.shapes.title and getattr(slide.shapes.title, "text", None):
            title = slide.shapes.title.text.strip()
            if title:
                slide_bits.append(f"Title: {title}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != title:
                    slide_bits.append(text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_bits.append("[Image on slide]")
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    slide_bits.append("Table:\n" + "\n".join(rows))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            slide_bits.append(f"Speaker notes: {notes}")
        if len(slide_bits) > 1:
            parts.append("\n".join(slide_bits))

    combined = "\n\n".join(parts).strip()
    if not combined:
        raise ValidationError(
            "EMPTY_EXTRACTION",
            "No readable text was found in this PowerPoint file.",
        )
    return ExtractionResult(text=combined, extractor="python-pptx")


def _extract_docx(path: Path) -> ExtractionResult:
    try:
        import docx
        from docx.enum.style import WD_STYLE_TYPE
    except ImportError as exc:
        raise _missing_dep(
            "EXTRACTION_UNAVAILABLE",
            "Word document processing is temporarily unavailable. Please try PDF or TXT.",
            cause=exc,
        ) from exc

    document = docx.Document(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
        if "heading" in style_name:
            parts.append(f"## {text}")
        elif paragraph.style and paragraph.style.type == WD_STYLE_TYPE.LIST:
            parts.append(f"- {text}")
        else:
            # Detect list-like paragraphs
            if text.startswith(("•", "-", "*", "–")):
                parts.append(text if text.startswith("-") else f"- {text.lstrip('•*-– ')}")
            else:
                parts.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"## Table {table_index}\n" + "\n".join(rows))

    combined = "\n\n".join(parts).strip()
    if not combined:
        raise ValidationError(
            "EMPTY_EXTRACTION",
            "No readable text was found in this Word document.",
        )
    return ExtractionResult(text=combined, extractor="python-docx")


def _extract_image_ocr(path: Path) -> ExtractionResult:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise _missing_dep(
            "EXTRACTION_UNAVAILABLE",
            "Image OCR is temporarily unavailable. Please upload a PDF, DOCX, PPTX, or TXT file instead.",
            cause=exc,
        ) from exc

    try:
        text = pytesseract.image_to_string(Image.open(path))
    except pytesseract.TesseractNotFoundError as exc:
        raise ValidationError(
            "OCR_UNAVAILABLE",
            "Image OCR is not configured on this server. Please upload a text-based PDF, DOCX, PPTX, or TXT file.",
        ) from exc
    except Exception as exc:
        raise ValidationError(
            "OCR_FAILED",
            "Could not read text from this image. Try a clearer scan or a different file format.",
        ) from exc

    if not text.strip():
        raise ValidationError(
            "EMPTY_EXTRACTION",
            "No text could be detected in this image. Try a higher-resolution image or a text document.",
        )
    return ExtractionResult(text=text.strip(), extractor="tesseract")


def _extract_video_whisper(path: Path) -> ExtractionResult:
    try:
        import whisper
    except ImportError as exc:
        raise _missing_dep(
            "EXTRACTION_UNAVAILABLE",
            "Video transcription is temporarily unavailable. Please upload a PDF, DOCX, PPTX, TXT, or image instead.",
            cause=exc,
        ) from exc

    try:
        model = whisper.load_model("base")
        result = model.transcribe(str(path))
        text = str(result.get("text") or "").strip()
    except Exception as exc:
        logger.exception("Whisper transcription failed")
        raise ValidationError(
            "TRANSCRIPTION_FAILED",
            "Could not transcribe this video. Ensure it is a valid MP4 with clear audio, or upload a text document instead.",
        ) from exc

    if not text:
        raise ValidationError(
            "EMPTY_EXTRACTION",
            "No speech could be transcribed from this video.",
        )
    return ExtractionResult(text=text, extractor="whisper")
